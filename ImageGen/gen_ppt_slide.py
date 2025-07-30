import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

dir_images = os.path.join(os.path.dirname(__file__), "ImageGen", "attrition_images") + os.sep

def generate_slides(image_pairs, comments):
    # Create a new PowerPoint presentation
    prs = Presentation()
    
	# Set slide size (16:9 aspect ratio)
    prs.slide_width = Inches(13.3333)
    prs.slide_height = Inches(7.5)
    
    # Generate slides for all image pairs
    for i in range(len(image_pairs)):
        # Add a blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Use blank layout
        # 1. Add left "Event 1" label
        left_box = slide.shapes.add_textbox(Inches(2.66), Inches(0.6), Inches(1.32), Inches(0.5))
        left_frame = left_box.text_frame

        left_frame.text = "Event 1"
        p = left_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Add left image (with error handling)
        try:
            slide.shapes.add_picture(
                dir_images + image_pairs[i][0], 
                Inches(0.93), Inches(1.36), 
                Inches(4.78), Inches(4.78)
            )
        except Exception as e:
            print(f"Failed to add left image: {e}")
        
        # 2. Add right "Event 2" label
        right_box = slide.shapes.add_textbox(Inches(9.35), Inches(0.6), Inches(1.32), Inches(0.5))
        right_frame = right_box.text_frame

        right_frame.text = "Event 2"
        p = right_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Add right image (with error handling)
        try:
            slide.shapes.add_picture(
                dir_images + image_pairs[i][1], 
                Inches(7.62), Inches(1.36), 
                Inches(4.78), Inches(4.78)
            )
        except Exception as e:
            print(f"Failed to add right image: {e}")
        
        # 3. Add "then" text in the center
        then_box = slide.shapes.add_textbox(Inches(6.23), Inches(3.13), Inches(0.88), Inches(0.5))
        then_frame = then_box.text_frame
        then_frame.text = "then"
        p = then_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
		
		# 4. Add footer text if comment is not "0"
        if comments[i] != "0":
            footer_box = slide.shapes.add_textbox(Inches(4.29), Inches(6.1), Inches(4.76), Inches(0.58))
            footer_frame = footer_box.text_frame
            footer_frame.text = comments[i]
            p = footer_frame.paragraphs[0]
            p.font.size = Pt(28)
            p.alignment = PP_ALIGN.CENTER
    
   	# Save the presentation
    prs.save("new_presentation.pptx")
    print("Slides generated and saved as new_presentation.pptx")

if __name__ == "__main__":
	# Paste the image pairs and comments here(the output of gen_ppt_data.py)
